import os
from datetime import datetime
from typing import Dict, Any, List, Optional
import gc
import matplotlib.pyplot as plt
import pandas as pd
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from docx import Document
from src.config import REPORTS_DIR

class ReportGenerator:
    """Class for generating reports in various formats."""
    
    def __init__(self):
        """Initialize the report generator."""
        self.reports_dir = REPORTS_DIR
        os.makedirs(self.reports_dir, exist_ok=True)
        # Set matplotlib to use a non-interactive backend with minimal memory usage
        plt.switch_backend('Agg')
        plt.rcParams['figure.max_open_warning'] = 0
        plt.rcParams['figure.dpi'] = 100
        plt.rcParams['figure.figsize'] = [10, 6]
    
    def generate_report(self, analysis_data: Dict[str, Any], report_type: str,
                       metrics: Optional[List[str]] = None,
                       branding: Optional[Dict[str, str]] = None) -> str:
        """Generate a report in the specified format."""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"report_{timestamp}"
            
            # Clean up any existing matplotlib figures
            plt.close('all')
            
            if report_type == 'pdf':
                return self._generate_pdf_report(analysis_data, filename, metrics, branding)
            elif report_type == 'pptx':
                return self._generate_pptx_report(analysis_data, filename, metrics, branding)
            elif report_type == 'docx':
                return self._generate_docx_report(analysis_data, filename, metrics, branding)
            elif report_type == 'excel':
                return self._generate_excel_report(analysis_data, filename, metrics)
            else:
                raise ValueError(f"Unsupported report type: {report_type}")
        finally:
            # Clean up memory
            gc.collect()
    
    def _generate_pdf_report(self, analysis_data: Dict[str, Any], filename: str,
                           metrics: Optional[List[str]] = None,
                           branding: Optional[Dict[str, str]] = None) -> str:
        """Generate a PDF report."""
        filepath = os.path.join(self.reports_dir, f"{filename}.pdf")
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []
        
        # Add title
        title = branding.get('company_name', 'ICP Analysis Report') if branding else 'ICP Analysis Report'
        elements.append(Paragraph(title, styles['Title']))
        elements.append(Spacer(1, 12))
        
        # Add overview section
        elements.append(Paragraph("Overview", styles['Heading1']))
        elements.append(Paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
        elements.append(Spacer(1, 12))
        
        # Add metrics table
        if metrics:
            data = [['Metric', 'Value']]
            for metric in metrics:
                if metric in analysis_data:
                    data.append([metric, str(analysis_data[metric])])
            
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 14),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 12),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            elements.append(table)
            elements.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(elements)
        return f"{filename}.pdf"
    
    def _generate_pptx_report(self, analysis_data: Dict[str, Any], filename: str,
                            metrics: Optional[List[str]] = None,
                            branding: Optional[Dict[str, str]] = None) -> str:
        """Generate a PowerPoint report."""
        filepath = os.path.join(self.reports_dir, f"{filename}.pptx")
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = branding.get('company_name', 'ICP Analysis Report') if branding else 'ICP Analysis Report'
        subtitle.text = f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Overview slide
        overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = overview_slide.shapes.title
        content = overview_slide.placeholders[1]
        
        title.text = "Overview"
        content.text = "Key metrics and insights from the analysis"
        
        # Metrics slide
        if metrics:
            metrics_slide = prs.slides.add_slide(prs.slide_layouts[2])
            title = metrics_slide.shapes.title
            content = metrics_slide.placeholders[1]
            
            title.text = "Key Metrics"
            content.text = "\n".join([f"{metric}: {analysis_data.get(metric, 'N/A')}" for metric in metrics])
        
        prs.save(filepath)
        return f"{filename}.pptx"
    
    def _generate_docx_report(self, analysis_data: Dict[str, Any], filename: str,
                            metrics: Optional[List[str]] = None,
                            branding: Optional[Dict[str, str]] = None) -> str:
        """Generate a Word report."""
        filepath = os.path.join(self.reports_dir, f"{filename}.docx")
        doc = Document()
        
        # Add title
        doc.add_heading(branding.get('company_name', 'ICP Analysis Report') if branding else 'ICP Analysis Report', 0)
        doc.add_paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph()
        
        # Add overview section
        doc.add_heading('Overview', level=1)
        doc.add_paragraph()
        
        # Add metrics table
        if metrics:
            doc.add_heading('Key Metrics', level=2)
            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid'
            
            # Add header row
            header_cells = table.rows[0].cells
            header_cells[0].text = 'Metric'
            header_cells[1].text = 'Value'
            
            # Add data rows
            for metric in metrics:
                if metric in analysis_data:
                    row_cells = table.add_row().cells
                    row_cells[0].text = metric
                    row_cells[1].text = str(analysis_data[metric])
            
            doc.add_paragraph()
        
        doc.save(filepath)
        return f"{filename}.docx"
    
    def _generate_excel_report(self, analysis_data: Dict[str, Any], filename: str,
                             metrics: Optional[List[str]] = None) -> str:
        """Generate an Excel report."""
        filepath = os.path.join(self.reports_dir, f"{filename}.xlsx")
        
        # Create DataFrame with only required columns and optimize memory usage
        if metrics:
            data = {metric: [analysis_data.get(metric, 'N/A')] for metric in metrics}
        else:
            # Only include non-dict and non-list values to reduce memory
            data = {key: [value] for key, value in analysis_data.items() 
                   if not isinstance(value, (dict, list))}
        
        # Create DataFrame with optimized memory usage
        df = pd.DataFrame(data)
        
        # Optimize DataFrame memory usage
        for col in df.columns:
            if df[col].dtype == 'object':
                df[col] = df[col].astype('category')
        
        # Save to Excel with optimized settings
        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Analysis Results')
        
        # Clear DataFrame from memory
        del df
        return f"{filename}.xlsx"
    
    def export_chart(self, chart_data: Dict[str, Any], chart_type: str,
                    title: str, filename: Optional[str] = None) -> str:
        """Export a chart as an image."""
        try:
            if not filename:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"chart_{timestamp}.png"
            
            filepath = os.path.join(self.reports_dir, filename)
            
            # Create figure with optimized settings
            plt.figure(figsize=(10, 6), dpi=100)
            
            # Plot data based on chart type
            if chart_type == 'bar':
                plt.bar(chart_data['labels'], chart_data['values'])
            elif chart_type == 'line':
                plt.plot(chart_data['labels'], chart_data['values'])
            elif chart_type == 'pie':
                plt.pie(chart_data['values'], labels=chart_data['labels'])
            elif chart_type == 'scatter':
                plt.scatter(chart_data['x'], chart_data['y'])
            
            # Customize chart
            plt.title(title)
            plt.xlabel(chart_data.get('xlabel', ''))
            plt.ylabel(chart_data.get('ylabel', ''))
            
            # Save chart with optimized settings
            plt.savefig(filepath, bbox_inches='tight', pad_inches=0.1, optimize=True)
            plt.close()
            
            return filename
        finally:
            # Clean up memory
            gc.collect() 